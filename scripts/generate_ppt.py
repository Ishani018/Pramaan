"""
Pramaan Feature Showcase – PowerPoint Generator
================================================
Generates a professional pitch deck for Project Pramaan.
Run: python scripts/generate_ppt.py
Output: Pramaan_Feature_Deck.pptx in project root
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── BRANDING ───────────────────────────────────────────────────────────
DARK_NAVY   = RGBColor(0x0A, 0x1F, 0x44)
BRAND_BLUE  = RGBColor(0x0F, 0x33, 0x6B)
MID_BLUE    = RGBColor(0x1E, 0x50, 0x8A)
LIGHT_BLUE  = RGBColor(0xE8, 0xEE, 0xF6)
ACCENT_RED  = RGBColor(0xCC, 0x22, 0x22)
ACCENT_GREEN= RGBColor(0x00, 0x7A, 0x4D)
AMBER       = RGBColor(0xCC, 0x7A, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF8, 0xF9, 0xFA)
GREY        = RGBColor(0x64, 0x74, 0x8B)
BLACK       = RGBColor(0x1A, 0x1A, 0x2E)

LOGO_PATH = os.path.join(os.path.dirname(__file__), "..", "backend", "app", "utils", "pramaan_logo.png")


def _add_bg(slide, color):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape_bg(slide, left, top, width, height, color, alpha=None):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_text_box(slide, left, top, width, height, text, font_size=18,
                  color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = font_name
    return txBox


def _add_bullet_list(slide, left, top, width, height, items, font_size=14,
                     color=BLACK, spacing=Pt(6)):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing

        # Handle (text, sub_color) tuples or plain strings
        if isinstance(item, tuple):
            text, item_color = item
        else:
            text, item_color = item, color

        r = p.add_run()
        r.text = f"▸  {text}"
        r.font.size = Pt(font_size)
        r.font.color.rgb = item_color
        r.font.name = "Calibri"

    return txBox


def _add_feature_card(slide, left, top, width, height, title, desc, icon_text="●", accent=BRAND_BLUE):
    """Add a feature card with icon, title, and description."""
    # Card background
    card = _add_shape_bg(slide, left, top, width, height, OFF_WHITE)

    # Accent bar on left
    _add_shape_bg(slide, left, top, Inches(0.06), height, accent)

    # Icon circle
    icon_shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left + Inches(0.2), top + Inches(0.2),
        Inches(0.4), Inches(0.4)
    )
    icon_shape.fill.solid()
    icon_shape.fill.fore_color.rgb = accent
    icon_shape.line.fill.background()
    tf = icon_shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = icon_text
    r.font.size = Pt(12)
    r.font.color.rgb = WHITE
    r.font.bold = True

    # Title
    _add_text_box(slide, left + Inches(0.75), top + Inches(0.12),
                  width - Inches(0.9), Inches(0.35),
                  title, font_size=13, color=DARK_NAVY, bold=True)

    # Description
    _add_text_box(slide, left + Inches(0.75), top + Inches(0.45),
                  width - Inches(0.9), height - Inches(0.55),
                  desc, font_size=10, color=GREY)


def _add_slide_number(slide, num, total):
    """Add slide number in bottom right."""
    _add_text_box(slide, Inches(8.8), Inches(7.1), Inches(1.0), Inches(0.3),
                  f"{num}/{total}", font_size=8, color=GREY, alignment=PP_ALIGN.RIGHT)


def _section_divider(prs, title, subtitle=""):
    """Create a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_bg(slide, DARK_NAVY)

    # Accent line
    _add_shape_bg(slide, Inches(1.0), Inches(3.2), Inches(1.5), Inches(0.04), ACCENT_RED)

    _add_text_box(slide, Inches(1.0), Inches(3.4), Inches(8.0), Inches(1.0),
                  title, font_size=32, color=WHITE, bold=True)

    if subtitle:
        _add_text_box(slide, Inches(1.0), Inches(4.3), Inches(8.0), Inches(0.8),
                      subtitle, font_size=16, color=LIGHT_BLUE)

    return slide


# ══════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    TOTAL_SLIDES = 16

    # ── SLIDE 1: TITLE ───────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, DARK_NAVY)

    # Logo
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(0.8), Inches(0.8), Inches(1.5))

    # Title
    _add_text_box(slide, Inches(0.8), Inches(2.8), Inches(8.5), Inches(1.2),
                  "PROJECT PRAMAAN", font_size=44, color=WHITE, bold=True)

    # Accent line
    _add_shape_bg(slide, Inches(0.8), Inches(4.0), Inches(2.0), Inches(0.05), ACCENT_RED)

    # Subtitle
    _add_text_box(slide, Inches(0.8), Inches(4.2), Inches(8.5), Inches(0.8),
                  "The Deterministic Credit Decisioning Engine", font_size=22, color=LIGHT_BLUE)

    _add_text_box(slide, Inches(0.8), Inches(5.0), Inches(8.5), Inches(0.6),
                  "Zero Hallucination  ·  33 Penalty Rules  ·  Full Explainability  ·  Institutional-Grade CAM",
                  font_size=13, color=GREY)

    _add_text_box(slide, Inches(0.8), Inches(6.2), Inches(8.5), Inches(0.5),
                  "Next-Generation Corporate Credit Appraisal for Indian Lending",
                  font_size=14, color=OFF_WHITE, bold=True)

    _add_slide_number(slide, 1, TOTAL_SLIDES)

    # ── SLIDE 2: PROBLEM STATEMENT ─────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_shape_bg(slide, Inches(0), Inches(0), Inches(10), Inches(1.2), DARK_NAVY)
    _add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                  "THE PROBLEM", font_size=28, color=WHITE, bold=True)

    problems = [
        ("Manual Credit Appraisal Takes 3–5 Days", "Credit officers manually cross-check annual reports, bank statements, GST filings, and external sources. Slow, error-prone, and inconsistent."),
        ("LLM-Based Solutions Hallucinate", "Generative AI models fabricate financial figures and compliance findings. Unacceptable risk for lending decisions worth crores."),
        ("Fraud Signals Are Buried in Data", "Circular trading, shell companies, GST mismatches, and silent restatements hide across 100+ page annual reports and thousands of bank transactions."),
        ("No Standardised Risk Framework", "Every credit officer applies different judgment. No consistent penalty system for rate adjustments and limit decisions."),
    ]

    for i, (title, desc) in enumerate(problems):
        y = Inches(1.6) + Inches(1.35) * i
        # Number circle
        num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.45), Inches(0.45))
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = ACCENT_RED
        num_shape.line.fill.background()
        tf = num_shape.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = str(i + 1)
        r.font.size = Pt(16)
        r.font.color.rgb = WHITE
        r.font.bold = True

        _add_text_box(slide, Inches(1.5), y - Inches(0.05), Inches(7.5), Inches(0.4),
                      title, font_size=15, color=DARK_NAVY, bold=True)
        _add_text_box(slide, Inches(1.5), y + Inches(0.35), Inches(7.5), Inches(0.6),
                      desc, font_size=11, color=GREY)

    _add_slide_number(slide, 2, TOTAL_SLIDES)

    # ── SLIDE 3: SOLUTION OVERVIEW ─────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_shape_bg(slide, Inches(0), Inches(0), Inches(10), Inches(1.2), DARK_NAVY)
    _add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                  "OUR SOLUTION: PRAMAAN", font_size=28, color=WHITE, bold=True)

    _add_text_box(slide, Inches(0.8), Inches(1.5), Inches(8.5), Inches(0.8),
                  "A deterministic credit engine that analyses annual reports, bank statements, and "
                  "external bureau data using zero-LLM architecture. Every finding is traceable. "
                  "Every decision is explainable.",
                  font_size=14, color=BLACK)

    pillars = [
        ("📄", "Upload & Analyze", "Annual report PDF + bank\nstatement CSV → instant\nanalysis pipeline"),
        ("🔍", "33 Penalty Rules", "Deterministic risk detection\nacross compliance, fraud,\nfinancials, operations"),
        ("🔗", "Cross-Verify", "Reconcile AR claims against\n8+ external sources: GST,\nMCA, eCourts, CIBIL, News"),
        ("📊", "Decision Engine", "Automated rate/limit\nadjustment with full\nwaterfall transparency"),
        ("📑", "CAM Export", "Institutional-grade Credit\nAppraisal Memo in 12\nbank-format sections"),
    ]

    for i, (icon, title, desc) in enumerate(pillars):
        x = Inches(0.4) + Inches(1.9) * i
        y = Inches(2.8)

        # Card
        _add_shape_bg(slide, x, y, Inches(1.7), Inches(3.5), LIGHT_BLUE)
        _add_shape_bg(slide, x, y, Inches(1.7), Inches(0.06), BRAND_BLUE)

        # Icon
        _add_text_box(slide, x, y + Inches(0.2), Inches(1.7), Inches(0.6),
                      icon, font_size=28, alignment=PP_ALIGN.CENTER)

        # Title
        _add_text_box(slide, x + Inches(0.1), y + Inches(0.85), Inches(1.5), Inches(0.5),
                      title, font_size=12, color=DARK_NAVY, bold=True, alignment=PP_ALIGN.CENTER)

        # Desc
        _add_text_box(slide, x + Inches(0.1), y + Inches(1.4), Inches(1.5), Inches(1.8),
                      desc, font_size=9, color=GREY, alignment=PP_ALIGN.CENTER)

    _add_slide_number(slide, 3, TOTAL_SLIDES)

    # ── SLIDE 4: ARCHITECTURE / PIPELINE ─────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_shape_bg(slide, Inches(0), Inches(0), Inches(10), Inches(1.2), DARK_NAVY)
    _add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                  "SYSTEM ARCHITECTURE", font_size=28, color=WHITE, bold=True)

    # Pipeline stages
    stages = [
        ("INPUT", "Annual Report PDF\nBank Statement CSV\nSite Visit Notes", BRAND_BLUE),
        ("DEEP READER", "PDF Extraction\nCompliance Scanner\nFinancial Extractor\nRating Extractor\nShareholding Scanner", MID_BLUE),
        ("EXTERNAL INTEL", "MCA21 Registry\neCourts Litigation\nNewsScanner Media\nPerfios GST\nKARZA Bureau", RGBColor(0x2D, 0x6A, 0x4F)),
        ("CROSS-VERIFY", "Claims Extraction\n8-Source Verification\nGST-Bank Recon\nCounterparty Intel\nNetwork Analysis", AMBER),
        ("DECISION", "33 Penalty Rules\nRate Waterfall\nLimit Adjustment\nCAM Generation\nFinal Recommendation", ACCENT_RED),
    ]

    for i, (title, items, color) in enumerate(stages):
        x = Inches(0.3) + Inches(1.95) * i
        y = Inches(1.8)

        # Stage box
        box = _add_shape_bg(slide, x, y, Inches(1.75), Inches(4.0), color)

        # Arrow between stages
        if i < len(stages) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, x + Inches(1.8), y + Inches(1.8),
                Inches(0.15), Inches(0.4)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GREY
            arrow.line.fill.background()

        # Title
        _add_text_box(slide, x + Inches(0.05), y + Inches(0.15),
                      Inches(1.65), Inches(0.4),
                      title, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        # Separator
        _add_shape_bg(slide, x + Inches(0.15), y + Inches(0.55), Inches(1.45), Inches(0.02), WHITE)

        # Items
        _add_text_box(slide, x + Inches(0.1), y + Inches(0.7),
                      Inches(1.55), Inches(3.0),
                      items, font_size=9, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Tech stack bar
    _add_shape_bg(slide, Inches(0.3), Inches(6.2), Inches(9.4), Inches(0.8), LIGHT_BLUE)
    _add_text_box(slide, Inches(0.5), Inches(6.25), Inches(9.0), Inches(0.3),
                  "TECH STACK", font_size=10, color=DARK_NAVY, bold=True)
    _add_text_box(slide, Inches(0.5), Inches(6.5), Inches(9.0), Inches(0.4),
                  "React + Vite + Tailwind  |  FastAPI + Uvicorn  |  PyMuPDF + Pandas  |  python-docx  |  "
                  "React Force Graph 3D  |  Recharts  |  MCA21 / eCourts / NewsAPI / Perfios / KARZA",
                  font_size=9, color=GREY)

    _add_slide_number(slide, 4, TOTAL_SLIDES)

    # ── SLIDE 5: SECTION DIVIDER — CORE FEATURES ─────────────────────
    slide = _section_divider(prs, "CORE FEATURES", "What makes Pramaan different")
    _add_slide_number(slide, 5, TOTAL_SLIDES)

    # ── SLIDE 6: GST-BANK RECONCILIATION ─────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_shape_bg(slide, Inches(0), Inches(0), Inches(10), Inches(1.2), DARK_NAVY)
    _add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                  "GST-BANK RECONCILIATION", font_size=28, color=WHITE, bold=True)

    _add_text_box(slide, Inches(0.8), Inches(1.5), Inches(8.5), Inches(0.7),
                  "Automatically reconciles declared GST turnover (Perfios) against actual bank "
                  "statement credits to detect revenue overstatement or diversion of funds.",
                  font_size=13, color=BLACK)

    # How it works
    _add_text_box(slide, Inches(0.8), Inches(2.4), Inches(4.0), Inches(0.4),
                  "How It Works", font_size=16, color=DARK_NAVY, bold=True)

    steps = [
        "Fetch GST turnover from Perfios API (GSTR-3B aggregate)",
        "Parse bank statement CSV → sum all credit entries",
        "Compute variance: (Bank Credits − GST Turnover) / GST Turnover",
        "Classify: ≤10% = MATCH | 10–20% = PARTIAL | >20% = MISMATCH",
        "MISMATCH triggers P-33: +125 bps rate penalty, −20% limit cut",
    ]
    _add_bullet_list(slide, Inches(0.8), Inches(2.9), Inches(4.5), Inches(3.0),
                     steps, font_size=11, color=BLACK)

    # Right side — result box
    _add_shape_bg(slide, Inches(5.8), Inches(2.4), Inches(3.8), Inches(4.0), LIGHT_BLUE)
    _add_shape_bg(slide, Inches(5.8), Inches(2.4), Inches(3.8), Inches(0.5), BRAND_BLUE)
    _add_text_box(slide, Inches(5.8), Inches(2.45), Inches(3.8), Inches(0.45),
                  "  Sample Output", font_size=13, color=WHITE, bold=True)

    sample_lines = [
        ("GST Turnover:", "₹0.58 Crore"),
        ("Bank Credits:", "₹0.57 Crore"),
        ("Variance:", "−1.7%"),
        ("Status:", "MATCH ✓"),
        ("", ""),
        ("GSTR 2A vs 3B:", "8.5% mismatch"),
        ("Filing Discipline:", "92.0/100"),
    ]
    y_pos = Inches(3.1)
    for label, value in sample_lines:
        if not label:
            y_pos += Inches(0.15)
            continue
        _add_text_box(slide, Inches(6.0), y_pos, Inches(1.6), Inches(0.3),
                      label, font_size=10, color=GREY, bold=True)
        val_color = ACCENT_GREEN if "MATCH" in value else (ACCENT_RED if "MISMATCH" in value else BLACK)
        _add_text_box(slide, Inches(7.6), y_pos, Inches(1.8), Inches(0.3),
                      value, font_size=10, color=val_color, bold=True)
        y_pos += Inches(0.32)

    # Rule tag
    _add_shape_bg(slide, Inches(0.8), Inches(6.3), Inches(8.8), Inches(0.6), RGBColor(0xFF, 0xF0, 0xF0))
    _add_text_box(slide, Inches(1.0), Inches(6.35), Inches(8.4), Inches(0.5),
                  "Rule P-33 (RECON-01): If variance exceeds 20%, rate increases by 125 bps and "
                  "limit is reduced by 20%. CA-certified reconciliation required as condition precedent.",
                  font_size=10, color=ACCENT_RED)

    _add_slide_number(slide, 6, TOTAL_SLIDES)

    # ── SLIDE 7: CIRCULAR TRADING & COUNTERPARTY INTEL ───────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_shape_bg(slide, Inches(0), Inches(0), Inches(10), Inches(1.2), DARK_NAVY)
    _add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                  "CIRCULAR TRADING & NETWORK INTELLIGENCE", font_size=26, color=WHITE, bold=True)

    # Left column
    _add_text_box(slide, Inches(0.8), Inches(1.5), Inches(4.5), Inches(0.4),
                  "Multi-Layer Fraud Detection", font_size=16, color=DARK_NAVY, bold=True)

    _add_text_box(slide, Inches(0.8), Inches(2.0), Inches(4.5), Inches(0.8),
                  "Pramaan detects circular trading through two independent channels, "
                  "then cross-references them for high-confidence fraud signals.",
                  font_size=11, color=BLACK)

    _add_feature_card(slide, Inches(0.8), Inches(3.0), Inches(4.2), Inches(1.1),
                      "Bank Statement Analysis [P-28]",
                      "Detects A→B and B→A flows within 7-day windows where both "
                      "amounts exceed ₹10,000. Flags round-trip fund routing.",
                      "₹", ACCENT_RED)

    _add_feature_card(slide, Inches(0.8), Inches(4.3), Inches(4.2), Inches(1.1),
                      "Counterparty Intelligence [P-06]",
                      "Network graph analysis: shared directors, common addresses, "
                      "shell company heuristics, circular fund flow chains.",
                      "🔗", BRAND_BLUE)

    _add_feature_card(slide, Inches(0.8), Inches(5.6), Inches(4.2), Inches(1.1),
                      "3D Network Visualization",
                      "Interactive force-directed graph showing entity relationships, "
                      "transaction volumes, and circular flow loops in real-time.",
                      "◉", MID_BLUE)

    # Right side — network viz description
    _add_shape_bg(slide, Inches(5.5), Inches(1.5), Inches(4.2), Inches(5.4), LIGHT_BLUE)
    _add_shape_bg(slide, Inches(5.5), Inches(1.5), Inches(4.2), Inches(0.5), BRAND_BLUE)
    _add_text_box(slide, Inches(5.5), Inches(1.55), Inches(4.2), Inches(0.45),
                  "  Network Graph Features", font_size=13, color=WHITE, bold=True)

    net_features = [
        "3D force-directed graph (React Force Graph)",
        "Nodes sized by transaction volume",
        "Red edges = circular flows detected",
        "Green nodes = MCA verified entities",
        "Orange nodes = shell company suspects",
        "Click to inspect counterparty details",
        "Shared director connections highlighted",
        "Auto-detects circular loop chains",
    ]
    _add_bullet_list(slide, Inches(5.7), Inches(2.2), Inches(3.8), Inches(3.0),
                     net_features, font_size=10, color=DARK_NAVY)

    # Detection stats
    _add_text_box(slide, Inches(5.7), Inches(5.3), Inches(3.8), Inches(0.3),
                  "Detection Capabilities:", font_size=11, color=DARK_NAVY, bold=True)

    det_items = [
        ("P-06 Circular Trading:", "+200 bps, −30% limit"),
        ("P-28 Bank Round-trips:", "+150 bps, −25% limit"),
        ("Shell Company Detection:", "MCA heuristics + address matching"),
    ]
    y = Inches(5.65)
    for label, val in det_items:
        _add_text_box(slide, Inches(5.7), y, Inches(1.8), Inches(0.25),
                      label, font_size=8, color=GREY, bold=True)
        _add_text_box(slide, Inches(7.5), y, Inches(2.0), Inches(0.25),
                      val, font_size=8, color=ACCENT_RED, bold=True)
        y += Inches(0.25)

    _add_slide_number(slide, 7, TOTAL_SLIDES)

    # ── SLIDE 8: SUPPLY CHAIN & COMPETITOR ANALYSIS ──────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_shape_bg(slide, Inches(0), Inches(0), Inches(10), Inches(1.2), DARK_NAVY)
    _add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                  "SUPPLY CHAIN & SECTOR BENCHMARKING", font_size=28, color=WHITE, bold=True)

    # Supply Chain Risk
    _add_text_box(slide, Inches(0.8), Inches(1.5), Inches(4.2), Inches(0.4),
                  "Supply Chain Risk Scoring", font_size=16, color=DARK_NAVY, bold=True)

    sc_items = [
        "Upstream supplier concentration & dependency risk",
        "Downstream buyer concentration & payment risk",
        "Weakest link identification in value chain",
        "Major supplier/buyer extraction from AR",
        "3-block framework: Supplier → Borrower → Buyer",
        "Risk bands: Low / Moderate / High with numeric scores",
    ]
    _add_bullet_list(slide, Inches(0.8), Inches(2.0), Inches(4.2), Inches(2.8),
                     sc_items, font_size=11, color=BLACK)

    # Sector Benchmark
    _add_text_box(slide, Inches(5.5), Inches(1.5), Inches(4.2), Inches(0.4),
                  "Sector Benchmark Analysis", font_size=16, color=DARK_NAVY, bold=True)

    bm_items = [
        "EBITDA margin vs sector average",
        "Revenue growth vs sector average",
        "Working capital cycle comparison",
        "Deviation scoring: OK / BELOW / CRITICAL",
        "Competitor structural advantages noted",
        "P-30: >25% underperformance → penalty",
    ]
    _add_bullet_list(slide, Inches(5.5), Inches(2.0), Inches(4.2), Inches(2.8),
                     bm_items, font_size=11, color=BLACK)

    # Competitor insight box
    _add_shape_bg(slide, Inches(0.8), Inches(5.0), Inches(8.8), Inches(2.0), LIGHT_BLUE)
    _add_shape_bg(slide, Inches(0.8), Inches(5.0), Inches(8.8), Inches(0.45), BRAND_BLUE)
    _add_text_box(slide, Inches(0.8), Inches(5.05), Inches(8.8), Inches(0.4),
                  "  Competitor Analysis in CAM", font_size=13, color=WHITE, bold=True)

    _add_text_box(slide, Inches(1.0), Inches(5.6), Inches(8.4), Inches(1.2),
                  "The CAM document automatically includes competitor analysis:\n"
                  "• Identifies structural advantages of competitors (e.g., vertical integration, captive sourcing)\n"
                  "• Generates deviation tables: Company vs Sector Benchmark with color-coded status\n"
                  "• Provides narrative on competitive positioning and market risk\n"
                  "• Flags CRITICAL underperformance that impacts credit assessment",
                  font_size=10, color=DARK_NAVY)

    _add_slide_number(slide, 8, TOTAL_SLIDES)

    # ── SLIDE 9: CROSS-VERIFICATION ENGINE ───────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_shape_bg(slide, Inches(0), Inches(0), Inches(10), Inches(1.2), DARK_NAVY)
    _add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                  "CROSS-VERIFICATION ENGINE", font_size=28, color=WHITE, bold=True)

    _add_text_box(slide, Inches(0.8), Inches(1.5), Inches(8.5), Inches(0.6),
                  "Extracts structured claims from the annual report and verifies each against "
                  "independent external data sources. Every claim gets a verdict.",
                  font_size=13, color=BLACK)

    # Claims verified
    claims = [
        ("Revenue Claim", "AR stated revenue vs Bank credits + GST turnover", "Bank + Perfios"),
        ("Profitability Claim", "AR stated margins vs Sector benchmark data", "Sector DB"),
        ("Compliance Claim", "AR statutory compliance vs CARO/auditor findings", "Compliance Scanner"),
        ("Going Concern", "AR going concern assertion vs Emphasis of Matter", "Audit Report"),
        ("Rating Claim", "AR stated credit rating vs External rating agencies", "Rating Extractor"),
        ("Litigation Status", "AR legal disclosures vs eCourts active cases", "eCourts API"),
        ("GST Reconciliation", "GST turnover (Perfios) vs Bank statement credits", "Bank CSV + Perfios"),
    ]

    # Table layout
    y_start = Inches(2.4)
    # Header
    _add_shape_bg(slide, Inches(0.8), y_start, Inches(8.5), Inches(0.4), DARK_NAVY)
    headers = [("Claim Type", 0.8, 2.0), ("Verification Logic", 2.8, 3.5), ("Data Source", 6.3, 2.0)]
    for text, x, w in headers:
        _add_text_box(slide, Inches(x), y_start + Inches(0.02), Inches(w), Inches(0.35),
                      f"  {text}", font_size=10, color=WHITE, bold=True)

    for i, (claim, logic, source) in enumerate(claims):
        y = y_start + Inches(0.4) + Inches(0.45) * i
        bg_color = LIGHT_BLUE if i % 2 == 0 else OFF_WHITE
        _add_shape_bg(slide, Inches(0.8), y, Inches(8.5), Inches(0.45), bg_color)
        _add_text_box(slide, Inches(0.9), y + Inches(0.02), Inches(1.8), Inches(0.4),
                      claim, font_size=10, color=DARK_NAVY, bold=True)
        _add_text_box(slide, Inches(2.9), y + Inches(0.02), Inches(3.3), Inches(0.4),
                      logic, font_size=9, color=BLACK)
        _add_text_box(slide, Inches(6.4), y + Inches(0.02), Inches(2.0), Inches(0.4),
                      source, font_size=9, color=GREY)

    # Verdicts
    _add_shape_bg(slide, Inches(0.8), Inches(5.8), Inches(8.5), Inches(0.8), LIGHT_BLUE)
    _add_text_box(slide, Inches(1.0), Inches(5.85), Inches(8.0), Inches(0.3),
                  "Possible Verdicts:", font_size=11, color=DARK_NAVY, bold=True)

    verdicts = [
        ("MATCH", ACCENT_GREEN, 1.0), ("PARTIAL MATCH", AMBER, 2.8),
        ("MISMATCH", ACCENT_RED, 5.0), ("UNVERIFIABLE", GREY, 7.0),
    ]
    for text, color, x in verdicts:
        _add_shape_bg(slide, Inches(x), Inches(6.2), Inches(1.5), Inches(0.3), color)
        _add_text_box(slide, Inches(x), Inches(6.2), Inches(1.5), Inches(0.3),
                      text, font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    _add_slide_number(slide, 9, TOTAL_SLIDES)

    # ── SLIDE 10: 33 PENALTY RULES ───────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_shape_bg(slide, Inches(0), Inches(0), Inches(10), Inches(1.2), DARK_NAVY)
    _add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                  "33 DETERMINISTIC PENALTY RULES", font_size=28, color=WHITE, bold=True)

    _add_text_box(slide, Inches(0.8), Inches(1.4), Inches(8.5), Inches(0.5),
                  "Every rule has a fixed rate penalty (bps) and limit reduction (%). "
                  "No black boxes. Full audit trail.",
                  font_size=12, color=GREY)

    categories = [
        ("COMPLIANCE\n& AUDIT", ["P-03: CARO Default (+150 bps)", "P-04: Emphasis of Matter (+75 bps)",
                                  "P-09: Silent Restatement (+200 bps)", "P-10: Auditor Change (+75 bps)"],
         BRAND_BLUE),
        ("FRAUD\nDETECTION", ["P-06: Circular Trading (+200 bps)", "P-28: Bank Round-trips (+150 bps)",
                               "P-29: Cash Spike near GST (+50 bps)", "P-05: Adverse Media (+50 bps)"],
         ACCENT_RED),
        ("FINANCIAL\nRISK", ["P-01: GST Revenue Gap (+100 bps)", "P-33: GST-Bank Mismatch (+125 bps)",
                              "P-30: Sector Underperform (+50 bps)", "P-31: Revenue Unverified (+100 bps)"],
         AMBER),
        ("GOVERNANCE\n& LEGAL", ["P-02: Director Network (+75 bps)", "P-13: Adverse Media (+50 bps)",
                                  "P-15: Court Proceedings (+100 bps)", "P-14: MCA Not Active (−100% limit)"],
         RGBColor(0x2D, 0x6A, 0x4F)),
        ("OPERATIONS\n& SITE", ["P-07: Site Visit Risk (+75 bps)", "P-19: Low Capacity (+100 bps)",
                                 "P-20: Premises Vacant (+150 bps)", "P-22: Mgmt Non-Coop (+125 bps)"],
         MID_BLUE),
    ]

    for i, (cat_title, rules, color) in enumerate(categories):
        x = Inches(0.3) + Inches(1.95) * i
        y = Inches(2.1)

        # Category header
        _add_shape_bg(slide, x, y, Inches(1.8), Inches(0.7), color)
        _add_text_box(slide, x, y + Inches(0.05), Inches(1.8), Inches(0.6),
                      cat_title, font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        # Rules
        for j, rule in enumerate(rules):
            ry = y + Inches(0.85) + Inches(0.65) * j
            _add_shape_bg(slide, x, ry, Inches(1.8), Inches(0.55), LIGHT_BLUE if j % 2 == 0 else OFF_WHITE)
            _add_text_box(slide, x + Inches(0.05), ry + Inches(0.02), Inches(1.7), Inches(0.5),
                          rule, font_size=8, color=DARK_NAVY)

    # Base rates
    _add_shape_bg(slide, Inches(0.3), Inches(5.0), Inches(9.4), Inches(0.5), DARK_NAVY)
    _add_text_box(slide, Inches(0.5), Inches(5.05), Inches(9.0), Inches(0.4),
                  "Base Rate: 9.00% p.a.  |  Base Limit: ₹10.00 Crore  |  "
                  "Rules stack additively  |  Manual review threshold: any HIGH severity rule",
                  font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    _add_slide_number(slide, 10, TOTAL_SLIDES)

    # ── SLIDE 11: CAM EXPORT ──────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_shape_bg(slide, Inches(0), Inches(0), Inches(10), Inches(1.2), DARK_NAVY)
    _add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                  "CREDIT APPRAISAL MEMO (CAM) EXPORT", font_size=26, color=WHITE, bold=True)

    _add_text_box(slide, Inches(0.8), Inches(1.5), Inches(8.5), Inches(0.6),
                  "One-click export of a professional .docx Credit Appraisal Memo styled to match "
                  "IDFC First Bank / HDFC Bank institutional CAM formats. 12 structured sections.",
                  font_size=13, color=BLACK)

    # Left — sections list
    cam_sections = [
        "1. Borrower Profile",
        "2. Proposal Summary & Recommendation",
        "3. Financial Analysis (Key Financials + Ratios)",
        "4. GST-Bank Reconciliation",
        "5. Counterparty Intelligence & Circular Trading",
        "6. Supply Chain & Competitor Analysis",
        "7. Verification Details (KYC, Legal, Media, Cross-Verify)",
        "8. Collateral & Security",
        "9. Site Visit / Field Verification",
        "10. Risk Rule Matrix & Penalty Waterfall",
        "11. Recommendation & Conditions Precedent",
        "12. Committee Decision (Signature Block)",
    ]
    _add_bullet_list(slide, Inches(0.8), Inches(2.3), Inches(4.5), Inches(4.5),
                     cam_sections, font_size=10, color=DARK_NAVY)

    # Right — features
    _add_shape_bg(slide, Inches(5.6), Inches(2.3), Inches(4.0), Inches(4.8), LIGHT_BLUE)
    _add_shape_bg(slide, Inches(5.6), Inches(2.3), Inches(4.0), Inches(0.5), BRAND_BLUE)
    _add_text_box(slide, Inches(5.6), Inches(2.35), Inches(4.0), Inches(0.45),
                  "  CAM Features", font_size=13, color=WHITE, bold=True)

    cam_features = [
        "Dark navy section header bars",
        "Alternating row shading on tables",
        "Color-coded status badges (RED/AMBER/GREEN)",
        "Pramaan logo in header",
        "CONFIDENTIAL watermark",
        "Competitor deviation table with status",
        "Rate & limit calculation waterfall",
        "Conditions precedent based on triggered rules",
        "Signature block for 3 authorities",
        "Full disclaimer & methodology note",
    ]
    _add_bullet_list(slide, Inches(5.8), Inches(3.0), Inches(3.6), Inches(3.8),
                     cam_features, font_size=9, color=DARK_NAVY, spacing=Pt(4))

    _add_slide_number(slide, 11, TOTAL_SLIDES)

    # ── SLIDE 12: DASHBOARD OVERVIEW ─────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_shape_bg(slide, Inches(0), Inches(0), Inches(10), Inches(1.2), DARK_NAVY)
    _add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                  "DASHBOARD & USER EXPERIENCE", font_size=28, color=WHITE, bold=True)

    tabs = [
        ("📊 Decision", "Penalty waterfall chart, final recommendation, rate/limit calculations, applied rules detail",
         BRAND_BLUE),
        ("🛡️ Compliance", "CARO 2020 findings, auditor qualifications, restatement detection, CEO ambidexterity scores",
         MID_BLUE),
        ("✓ Verify", "GST reconciliation hero card, supply chain risk panel, network analysis, claim verification cards",
         ACCENT_GREEN),
        ("🔍 Intelligence", "Bank statement analysis, adverse media scanner, eCourts litigation, BSE annual report search",
         AMBER),
        ("📈 Financials", "Sector benchmarks, restatement analysis, financial trends, MD&A sentiment analysis",
         RGBColor(0x6B, 0x21, 0xA8)),
    ]

    for i, (tab_name, tab_desc, color) in enumerate(tabs):
        y = Inches(1.5) + Inches(1.1) * i
        _add_shape_bg(slide, Inches(0.8), y, Inches(8.5), Inches(0.95), LIGHT_BLUE if i % 2 == 0 else OFF_WHITE)
        _add_shape_bg(slide, Inches(0.8), y, Inches(0.08), Inches(0.95), color)

        _add_text_box(slide, Inches(1.1), y + Inches(0.05), Inches(2.5), Inches(0.35),
                      tab_name, font_size=14, color=DARK_NAVY, bold=True)
        _add_text_box(slide, Inches(1.1), y + Inches(0.4), Inches(8.0), Inches(0.5),
                      tab_desc, font_size=10, color=GREY)

    # Design note
    _add_shape_bg(slide, Inches(0.8), Inches(7.0) - Inches(0.5), Inches(8.5), Inches(0.5), DARK_NAVY)
    _add_text_box(slide, Inches(1.0), Inches(7.0) - Inches(0.45), Inches(8.0), Inches(0.4),
                  "Brutalist cream/ink/red design system  ·  Fully responsive  ·  Real-time analysis  ·  One-click CAM export",
                  font_size=10, color=WHITE, alignment=PP_ALIGN.CENTER)

    _add_slide_number(slide, 12, TOTAL_SLIDES)

    # ── SLIDE 13: EXTERNAL DATA SOURCES ──────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_shape_bg(slide, Inches(0), Inches(0), Inches(10), Inches(1.2), DARK_NAVY)
    _add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                  "EXTERNAL DATA INTEGRATION", font_size=28, color=WHITE, bold=True)

    sources = [
        ("MCA21 Registry", "Company status, CIN, directors,\ncharges, paid-up capital", "Live API", BRAND_BLUE),
        ("Perfios GST", "GSTR-3B turnover, 2A vs 3B\nmismatch, filing discipline", "Mock (Demo)", MID_BLUE),
        ("KARZA Bureau", "CIBIL score, litigation history,\ncharge holders, DPD analysis", "Mock (Demo)", RGBColor(0x2D, 0x6A, 0x4F)),
        ("eCourts API", "Active cases by NCLT, DRT,\nHigh Court; risk classification", "Live API", ACCENT_RED),
        ("NewsAPI", "Real-time adverse media scan:\nfraud, raid, SEBI, insolvency", "Live API", AMBER),
        ("BSE Database", "Annual report downloads from\nBSE corporate filings", "Live API", RGBColor(0x6B, 0x21, 0xA8)),
    ]

    for i, (name, desc, status, color) in enumerate(sources):
        col = i % 3
        row = i // 3
        x = Inches(0.5) + Inches(3.1) * col
        y = Inches(1.6) + Inches(2.8) * row

        # Card
        _add_shape_bg(slide, x, y, Inches(2.8), Inches(2.3), LIGHT_BLUE)
        _add_shape_bg(slide, x, y, Inches(2.8), Inches(0.06), color)

        # Name
        _add_text_box(slide, x + Inches(0.15), y + Inches(0.2), Inches(2.5), Inches(0.35),
                      name, font_size=14, color=DARK_NAVY, bold=True)

        # Description
        _add_text_box(slide, x + Inches(0.15), y + Inches(0.6), Inches(2.5), Inches(1.0),
                      desc, font_size=10, color=BLACK)

        # Status badge
        badge_color = ACCENT_GREEN if "Live" in status else AMBER
        _add_shape_bg(slide, x + Inches(0.15), y + Inches(1.7), Inches(1.2), Inches(0.3), badge_color)
        _add_text_box(slide, x + Inches(0.15), y + Inches(1.7), Inches(1.2), Inches(0.3),
                      status, font_size=8, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    _add_slide_number(slide, 13, TOTAL_SLIDES)

    # ── SLIDE 14: DIFFERENTIATORS ────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_shape_bg(slide, Inches(0), Inches(0), Inches(10), Inches(1.2), DARK_NAVY)
    _add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                  "WHY PRAMAAN?", font_size=28, color=WHITE, bold=True)

    diffs = [
        ("Zero Hallucination", "Pure keyword/regex/pattern matching. No generative AI. "
         "Every finding is traceable to the source document, page, and line.",
         "Unlike LLM-based credit tools, Pramaan never fabricates financial figures.", "0%", BRAND_BLUE),
        ("Full Explainability", "33 deterministic rules with fixed penalties. Every rate adjustment "
         "and limit reduction links to a specific rule and evidence.",
         "Credit officers can audit the entire decision chain in seconds.", "33", ACCENT_RED),
        ("Indian-Centric", "Built for CARO 2020, MCA21, GSTR-2A/3B, eCourts India, "
         "Indian sector benchmarks, and Loughran-McDonald financial lexicon.",
         "Not a Western tool adapted for India — built from scratch for Indian lending.", "🇮🇳", ACCENT_GREEN),
        ("Multi-Source", "Cross-verifies annual report claims against 8+ independent sources: "
         "bank statements, GST, CIBIL, MCA, eCourts, news, sector benchmarks.",
         "Single-source analysis is unreliable. Pramaan triangulates across all sources.", "8+", AMBER),
    ]

    for i, (title, desc, why, badge, color) in enumerate(diffs):
        y = Inches(1.5) + Inches(1.4) * i

        # Badge circle
        badge_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.1),
            Inches(0.6), Inches(0.6)
        )
        badge_shape.fill.solid()
        badge_shape.fill.fore_color.rgb = color
        badge_shape.line.fill.background()
        tf = badge_shape.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = badge
        r.font.size = Pt(14)
        r.font.color.rgb = WHITE
        r.font.bold = True

        _add_text_box(slide, Inches(1.7), y, Inches(7.5), Inches(0.35),
                      title, font_size=15, color=DARK_NAVY, bold=True)
        _add_text_box(slide, Inches(1.7), y + Inches(0.35), Inches(7.5), Inches(0.4),
                      desc, font_size=10, color=BLACK)
        _add_text_box(slide, Inches(1.7), y + Inches(0.75), Inches(7.5), Inches(0.3),
                      why, font_size=9, color=GREY, bold=True)

    _add_slide_number(slide, 14, TOTAL_SLIDES)

    # ── SLIDE 15: DEMO FLOW ──────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_shape_bg(slide, Inches(0), Inches(0), Inches(10), Inches(1.2), DARK_NAVY)
    _add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                  "LIVE DEMO FLOW", font_size=28, color=WHITE, bold=True)

    demo_steps = [
        ("STEP 1", "Upload Annual Report PDF", "Drag & drop any Indian company's annual report. "
         "The Deep Reader pipeline extracts financials, compliance, ratings, shareholding, and MD&A in seconds."),
        ("STEP 2", "Upload Bank Statement CSV", "Optional: upload a bank statement CSV to enable "
         "GST reconciliation, circular transaction detection, and counterparty intelligence."),
        ("STEP 3", "Review Dashboard", "Navigate 5 tabs: Decision waterfall, Compliance findings, "
         "Cross-verification results, Network intelligence, Financial analysis."),
        ("STEP 4", "Inspect Findings", "Click into GST reconciliation, circular trading network, "
         "supply chain risk, sector benchmarks. Every finding links to evidence."),
        ("STEP 5", "Export CAM", "One click → professional 12-section Word document. "
         "Ready for committee review with signature blocks."),
    ]

    for i, (step, title, desc) in enumerate(demo_steps):
        y = Inches(1.5) + Inches(1.1) * i
        bg = LIGHT_BLUE if i % 2 == 0 else OFF_WHITE
        _add_shape_bg(slide, Inches(0.8), y, Inches(8.5), Inches(0.95), bg)

        # Step badge
        _add_shape_bg(slide, Inches(0.8), y, Inches(0.9), Inches(0.95), BRAND_BLUE)
        _add_text_box(slide, Inches(0.8), y + Inches(0.1), Inches(0.9), Inches(0.4),
                      step, font_size=8, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        _add_text_box(slide, Inches(1.9), y + Inches(0.05), Inches(7.2), Inches(0.35),
                      title, font_size=13, color=DARK_NAVY, bold=True)
        _add_text_box(slide, Inches(1.9), y + Inches(0.4), Inches(7.2), Inches(0.5),
                      desc, font_size=10, color=GREY)

    _add_slide_number(slide, 15, TOTAL_SLIDES)

    # ── SLIDE 16: CLOSING / CONTACT ──────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, DARK_NAVY)

    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(4.0), Inches(1.0), Inches(2.0))

    _add_text_box(slide, Inches(1.0), Inches(3.3), Inches(8.0), Inches(0.8),
                  "PROJECT PRAMAAN", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    _add_shape_bg(slide, Inches(4.2), Inches(4.1), Inches(1.6), Inches(0.04), ACCENT_RED)

    _add_text_box(slide, Inches(1.0), Inches(4.3), Inches(8.0), Inches(0.6),
                  "The Deterministic Credit Decisioning Engine", font_size=18, color=LIGHT_BLUE,
                  alignment=PP_ALIGN.CENTER)

    # Stats
    stats = [
        ("33", "Penalty Rules"),
        ("8+", "Data Sources"),
        ("12", "CAM Sections"),
        ("0", "Hallucinations"),
    ]
    for i, (num, label) in enumerate(stats):
        x = Inches(1.5) + Inches(2.0) * i
        _add_text_box(slide, x, Inches(5.3), Inches(1.5), Inches(0.6),
                      num, font_size=32, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        _add_text_box(slide, x, Inches(5.8), Inches(1.5), Inches(0.4),
                      label, font_size=11, color=GREY, alignment=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(1.0), Inches(6.6), Inches(8.0), Inches(0.4),
                  "Zero LLM  ·  Full Explainability  ·  Indian-Centric  ·  Institutional Grade",
                  font_size=12, color=GREY, alignment=PP_ALIGN.CENTER)

    _add_slide_number(slide, 16, TOTAL_SLIDES)

    # ── SAVE ──────────────────────────────────────────────────────────
    output_path = os.path.join(os.path.dirname(__file__), "..", "Pramaan_Deck.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {os.path.abspath(output_path)}")
    print(f"Total slides: {TOTAL_SLIDES}")
    return output_path


if __name__ == "__main__":
    build_presentation()
